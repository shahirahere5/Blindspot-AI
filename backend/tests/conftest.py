"""Shared pytest fixtures: an isolated TestClient and generated sample files."""

from __future__ import annotations

import io
import shutil
from pathlib import Path

import pytest
from docx import Document as DocxDocument
from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture()
def client(tmp_path, monkeypatch):
    """
    Provide a TestClient wired to isolated, temporary storage directories so
    tests never touch (or pollute) the real backend/data folder.
    """
    uploads_dir = tmp_path / "uploads"
    documents_dir = tmp_path / "documents"
    vector_store_dir = tmp_path / "vector_store"
    version_groups_dir = tmp_path / "version_groups"
    comparison_cache_dir = tmp_path / "comparison_cache"
    graph_store_dir = tmp_path / "knowledge_graph"
    conversation_store_dir = tmp_path / "conversations"
    uploads_dir.mkdir()
    documents_dir.mkdir()
    vector_store_dir.mkdir()
    version_groups_dir.mkdir()
    comparison_cache_dir.mkdir()
    graph_store_dir.mkdir()
    conversation_store_dir.mkdir()

    import config

    monkeypatch.setattr(config, "UPLOADS_DIR", uploads_dir)
    monkeypatch.setattr(config, "DOCUMENTS_DIR", documents_dir)
    monkeypatch.setattr(config, "RAG_VECTOR_STORE_PATH", vector_store_dir)
    monkeypatch.setattr(config, "VERSION_GROUPS_DIR", version_groups_dir)
    monkeypatch.setattr(config, "COMPARISON_CACHE_DIR", comparison_cache_dir)
    monkeypatch.setattr(config, "GRAPH_STORE_DIR", graph_store_dir)
    monkeypatch.setattr(config, "CONVERSATION_STORE_DIR", conversation_store_dir)

    import storage.document_store as store_module

    monkeypatch.setattr(
        store_module,
        "document_store",
        store_module.DocumentStore(uploads_dir, documents_dir),
    )
    # api.documents imports `document_store` directly, so patch it there too.
    import api.documents as documents_module

    monkeypatch.setattr(documents_module, "document_store", store_module.document_store)

    import storage.version_store as version_store_module

    isolated_version_store = version_store_module.VersionStore(version_groups_dir)
    monkeypatch.setattr(version_store_module, "version_store", isolated_version_store)
    monkeypatch.setattr(documents_module, "version_store", isolated_version_store)

    # services.document_service (Phase 2) also imports `document_store`
    # directly, so it needs the same patch to see the isolated test storage.
    import services.document_service as document_service_module

    monkeypatch.setattr(
        document_service_module, "document_store", store_module.document_store
    )

    # Phase 4: isolate the local vector store the same way, and patch every
    # module that imports the singleton directly (services.retrieval_service).
    import storage.vector_store as vector_store_module

    isolated_vector_store = vector_store_module.SimpleVectorStore(vector_store_dir)
    monkeypatch.setattr(vector_store_module, "vector_store", isolated_vector_store)

    import services.retrieval_service as retrieval_service_module

    monkeypatch.setattr(
        retrieval_service_module, "vector_store", isolated_vector_store
    )

    import storage.comparison_cache as comparison_cache_module

    isolated_comparison_cache = comparison_cache_module.ComparisonCache(comparison_cache_dir)
    monkeypatch.setattr(comparison_cache_module, "comparison_cache", isolated_comparison_cache)

    import services.comparison_service as comparison_service_module

    monkeypatch.setattr(comparison_service_module, "version_store", isolated_version_store)
    monkeypatch.setattr(comparison_service_module, "comparison_cache", isolated_comparison_cache)

    import storage.graph_store as graph_store_module

    isolated_graph_store = graph_store_module.GraphStore(graph_store_dir)
    monkeypatch.setattr(graph_store_module, "graph_store", isolated_graph_store)

    import services.graph_ingestion_service as graph_ingestion_module
    import services.graph_service as graph_service_module

    monkeypatch.setattr(graph_ingestion_module, "graph_store", isolated_graph_store)
    monkeypatch.setattr(graph_ingestion_module, "version_store", isolated_version_store)
    monkeypatch.setattr(graph_service_module, "graph_store", isolated_graph_store)
    monkeypatch.setattr(graph_service_module, "version_store", isolated_version_store)

    import storage.conversation_store as conversation_store_module
    isolated_conversation_store = conversation_store_module.ConversationStore(conversation_store_dir)
    monkeypatch.setattr(conversation_store_module, "conversation_store", isolated_conversation_store)

    import services.conversation_service as conversation_service_module
    monkeypatch.setattr(conversation_service_module, "conversation_store", isolated_conversation_store)
    monkeypatch.setattr(conversation_service_module, "version_store", isolated_version_store)

    import main as main_module

    with TestClient(main_module.app) as test_client:
        yield test_client


@pytest.fixture()
def sample_txt_bytes() -> bytes:
    return "Hello Blind Spot AI.\nThis is a plain text test file.".encode("utf-8")


@pytest.fixture()
def sample_pdf_bytes() -> bytes:
    import fitz

    pdf = fitz.open()
    for i in range(3):
        page = pdf.new_page()
        page.insert_text((72, 72), f"This is page {i + 1} of the test PDF document.")
    buffer = io.BytesIO()
    pdf.save(buffer)
    pdf.close()
    return buffer.getvalue()


@pytest.fixture()
def sample_scanned_pdf_bytes() -> bytes:
    """A PDF with pages but no text layer, simulating a scanned document."""
    import fitz

    pdf = fitz.open()
    pdf.new_page()
    pdf.new_page()
    buffer = io.BytesIO()
    pdf.save(buffer)
    pdf.close()
    return buffer.getvalue()


@pytest.fixture()
def sample_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_heading("Test Document", level=1)
    document.add_paragraph("This is the first paragraph.")
    document.add_paragraph("This is the second paragraph.")

    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value 1"
    table.cell(1, 1).text = "Value 2"

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


@pytest.fixture()
def sample_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Our Pitch"
    body = slide.placeholders[1]
    body.text_frame.text = "We solve a real problem."

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "Market"
    slide2.placeholders[1].text_frame.text = "The market is huge."

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.fixture()
def sample_png_bytes() -> bytes:
    image = Image.new("RGB", (1920, 1080), color=(120, 200, 80))
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


@pytest.fixture()
def uploaded_txt_document_id(client, sample_txt_bytes) -> str:
    """Upload a small TXT document and return its document_id."""
    response = client.post(
        "/api/documents/upload",
        files={"file": ("pitch.txt", sample_txt_bytes, "text/plain")},
    )
    assert response.status_code == 200
    return response.json()["document_id"]


@pytest.fixture()
def uploaded_png_document_id(client, sample_png_bytes) -> str:
    """Upload a PNG (pending multimodal analysis) and return its document_id."""
    response = client.post(
        "/api/documents/upload",
        files={"file": ("photo.png", sample_png_bytes, "image/png")},
    )
    assert response.status_code == 200
    return response.json()["document_id"]

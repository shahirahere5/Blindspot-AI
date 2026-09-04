"""Visual candidate selection, normalization, provenance, and degradation tests."""

from __future__ import annotations

import io

import fitz
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

import config
import services.multimodal_service as multimodal_service
from ai.vision.base import VisionClient, VisionConfigurationError, VisionTimeoutError
from processing.image import ImageProcessor
from processing.pdf import PDFProcessor
from processing.pptx import PPTXProcessor
from schemas.document import DocumentStatus
from schemas.vision import VisualAnalysis


class FakeVisionClient(VisionClient):
    def __init__(self, fail_locations: set[int] | None = None) -> None:
        self.fail_locations = fail_locations or set()
        self.calls: list[tuple[str, int]] = []

    @property
    def provider_name(self) -> str:
        return "fake-vision"

    @property
    def model_name(self) -> str:
        return "fake-model"

    async def analyze_image(self, image_bytes, mime_type, source_type, source_location):
        self.calls.append((source_type, source_location))
        if source_location in self.fail_locations:
            raise VisionTimeoutError("raw provider timeout details")
        return VisualAnalysis(
            visible_text="Visible evidence in the supplied image",
            summary="Visual summary for the supplied evidence",
            key_evidence=["A decision-relevant chart trend"],
        )


@pytest.mark.asyncio
async def test_standalone_image_becomes_analyzable_visual_evidence(
    tmp_path, monkeypatch, sample_png_bytes
):
    monkeypatch.setattr(config, "MULTIMODAL_ENABLED", True)
    fake = FakeVisionClient()
    monkeypatch.setattr(multimodal_service, "get_vision_client", lambda: fake)
    path = tmp_path / "evidence.png"
    path.write_bytes(sample_png_bytes)
    document = ImageProcessor().process(path, "doc_image", "evidence.png")

    result = await multimodal_service.enrich_document_with_visuals(document, path)

    assert result.status == DocumentStatus.PROCESSED
    assert result.content[0].type.value == "image"
    assert result.content[0].location == 1
    assert "Visual summary" in result.content[0].text
    assert result.content[0].extra["visual_source"] == "multimodal"
    assert result.metadata["multimodal"]["status"] == "completed"


@pytest.mark.asyncio
async def test_text_pdf_does_not_construct_or_call_vision(
    tmp_path, monkeypatch, sample_pdf_bytes
):
    monkeypatch.setattr(config, "MULTIMODAL_ENABLED", True)
    monkeypatch.setattr(
        multimodal_service,
        "get_vision_client",
        lambda: (_ for _ in ()).throw(AssertionError("vision should not be used")),
    )
    path = tmp_path / "text.pdf"
    path.write_bytes(sample_pdf_bytes)
    document = PDFProcessor().process(path, "doc_text", "text.pdf")

    result = await multimodal_service.enrich_document_with_visuals(document, path)

    assert result.status == DocumentStatus.PROCESSED
    assert "multimodal" not in result.metadata


@pytest.mark.asyncio
async def test_scanned_pdf_pages_gain_visual_text_and_keep_page_locations(
    tmp_path, monkeypatch, sample_scanned_pdf_bytes
):
    monkeypatch.setattr(config, "MULTIMODAL_ENABLED", True)
    fake = FakeVisionClient()
    monkeypatch.setattr(multimodal_service, "get_vision_client", lambda: fake)
    path = tmp_path / "scan.pdf"
    path.write_bytes(sample_scanned_pdf_bytes)
    document = PDFProcessor().process(path, "doc_scan", "scan.pdf")

    result = await multimodal_service.enrich_document_with_visuals(document, path)

    assert result.status == DocumentStatus.PROCESSED
    assert [block.location for block in result.content] == [1, 2]
    assert all("Visual summary" in block.text for block in result.content)
    # The two empty rendered pages are byte-identical, so the provider result
    # is safely reused while the server still assigns both page locations.
    assert fake.calls == [("page", 1)]
    assert result.metadata["multimodal"]["attempted"] == 2


@pytest.mark.asyncio
async def test_hybrid_pdf_keeps_text_and_adds_visual_observations(
    tmp_path, monkeypatch
):
    monkeypatch.setattr(config, "MULTIMODAL_ENABLED", True)
    fake = FakeVisionClient()
    monkeypatch.setattr(multimodal_service, "get_vision_client", lambda: fake)

    image_buffer = io.BytesIO()
    Image.new("RGB", (800, 600), "blue").save(image_buffer, format="PNG")
    path = tmp_path / "hybrid.pdf"
    pdf = fitz.open()
    page = pdf.new_page(width=600, height=800)
    page.insert_text((40, 40), "Extracted strategy text remains available for review.")
    page.insert_image(fitz.Rect(40, 100, 560, 700), stream=image_buffer.getvalue())
    pdf.save(path)
    pdf.close()
    document = PDFProcessor().process(path, "doc_hybrid", "hybrid.pdf")

    result = await multimodal_service.enrich_document_with_visuals(document, path)

    assert "Extracted strategy text remains" in result.content[0].text
    assert "[Visual observations]" in result.content[0].text
    assert fake.calls == [("page", 1)]


@pytest.mark.asyncio
async def test_partial_pdf_failure_preserves_successful_content_and_safe_warning(
    tmp_path, monkeypatch
):
    monkeypatch.setattr(config, "MULTIMODAL_ENABLED", True)
    fake = FakeVisionClient(fail_locations={2})
    monkeypatch.setattr(multimodal_service, "get_vision_client", lambda: fake)
    path = tmp_path / "partial.pdf"
    pdf = fitz.open()
    pdf.new_page()
    second_page = pdf.new_page()
    second_page.draw_rect(fitz.Rect(40, 40, 200, 200), color=(1, 0, 0), fill=(1, 0, 0))
    pdf.save(path)
    pdf.close()
    document = PDFProcessor().process(path, "doc_partial", "partial.pdf")

    result = await multimodal_service.enrich_document_with_visuals(document, path)

    assert result.status == DocumentStatus.PROCESSED
    assert result.metadata["multimodal"]["status"] == "partial"
    assert any("Page 2" in warning for warning in result.warnings)
    assert all("raw provider" not in warning for warning in result.warnings)


@pytest.mark.asyncio
async def test_missing_provider_configuration_degrades_without_internal_details(
    tmp_path, monkeypatch, sample_png_bytes
):
    monkeypatch.setattr(config, "MULTIMODAL_ENABLED", True)
    monkeypatch.setattr(
        multimodal_service,
        "get_vision_client",
        lambda: (_ for _ in ()).throw(
            VisionConfigurationError("VISION_API_KEY missing in private .env")
        ),
    )
    path = tmp_path / "image.png"
    path.write_bytes(sample_png_bytes)
    document = ImageProcessor().process(path, "doc_unavailable", "image.png")

    result = await multimodal_service.enrich_document_with_visuals(document, path)

    assert result.status == DocumentStatus.PENDING_MULTIMODAL_ANALYSIS
    assert "Visual analysis is currently unavailable." in result.warnings
    assert all("VISION_API_KEY" not in warning for warning in result.warnings)


@pytest.mark.asyncio
async def test_pptx_picture_keeps_slide_provenance(
    tmp_path, monkeypatch, sample_png_bytes
):
    monkeypatch.setattr(config, "MULTIMODAL_ENABLED", True)
    fake = FakeVisionClient()
    monkeypatch.setattr(multimodal_service, "get_vision_client", lambda: fake)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(sample_png_bytes), Inches(1), Inches(1), width=Inches(5))
    path = tmp_path / "visual.pptx"
    presentation.save(path)
    document = PPTXProcessor().process(path, "doc_slides", "visual.pptx")

    result = await multimodal_service.enrich_document_with_visuals(document, path)

    assert result.content[0].location == 1
    assert "Visual summary for the supplied evidence" in result.content[0].text
    assert fake.calls == [("slide", 1)]


def test_image_pixel_limit_rejects_decompression_risk(
    tmp_path, monkeypatch, sample_png_bytes
):
    from processing.exceptions import DocumentProcessingError

    monkeypatch.setattr(config, "VISION_MAX_IMAGE_PIXELS", 10)
    path = tmp_path / "large.png"
    path.write_bytes(sample_png_bytes)

    with pytest.raises(DocumentProcessingError, match="pixel limit"):
        ImageProcessor().process(path, "doc_large", "large.png")


@pytest.mark.asyncio
async def test_visual_item_limit_caps_provider_calls(tmp_path, monkeypatch):
    monkeypatch.setattr(config, "MULTIMODAL_ENABLED", True)
    monkeypatch.setattr(config, "VISION_MAX_ITEMS_PER_DOCUMENT", 2)
    fake = FakeVisionClient()
    monkeypatch.setattr(multimodal_service, "get_vision_client", lambda: fake)
    path = tmp_path / "many-pages.pdf"
    pdf = fitz.open()
    for color in ((1, 0, 0), (0, 1, 0), (0, 0, 1)):
        page = pdf.new_page()
        page.draw_rect(fitz.Rect(30, 30, 180, 180), color=color, fill=color)
    pdf.save(path)
    pdf.close()
    document = PDFProcessor().process(path, "doc_many", "many-pages.pdf")

    result = await multimodal_service.enrich_document_with_visuals(document, path)

    assert len(fake.calls) == 2
    assert result.metadata["multimodal"]["skipped"] == 1
    assert any("limited to 2 items" in warning for warning in result.warnings)

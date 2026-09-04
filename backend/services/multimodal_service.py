"""Optional visual-evidence enrichment for normalized documents."""

from __future__ import annotations

import asyncio
import hashlib
import logging
import math
from dataclasses import dataclass
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import config
from ai.vision.base import VisionClient, VisionError, VisionResponseError
from ai.vision.factory import get_vision_client
from schemas.document import (
    ContentBlock,
    ContentBlockType,
    DocumentStatus,
    FileType,
    NormalizedDocument,
)
from schemas.vision import VisualAnalysis

logger = logging.getLogger("blindspot.multimodal")


@dataclass(frozen=True)
class VisualCandidate:
    image_bytes: bytes
    mime_type: str
    source_type: ContentBlockType
    location: int


def _safe_warning(source_type: ContentBlockType, location: int) -> str:
    label = source_type.value.capitalize()
    return f"Visual analysis was unavailable for {label} {location}."


def _append_once(values: list[str], value: str) -> None:
    if value not in values:
        values.append(value)


def _page_image_coverage(page: fitz.Page) -> float:
    page_area = max(1.0, page.rect.width * page.rect.height)
    covered_area = 0.0
    for image in page.get_images(full=True):
        try:
            for rect in page.get_image_rects(image[0]):
                clipped = rect & page.rect
                if not clipped.is_empty:
                    covered_area += clipped.width * clipped.height
        except Exception:  # malformed image metadata must not fail text extraction
            continue
    return min(1.0, covered_area / page_area)


def _render_pdf_page(page: fitz.Page) -> bytes:
    scale = max(0.5, config.VISION_PDF_RENDER_SCALE)
    pixel_count = page.rect.width * scale * page.rect.height * scale
    max_pixels = max(1, config.VISION_MAX_IMAGE_PIXELS)
    if pixel_count > max_pixels:
        scale *= math.sqrt(max_pixels / pixel_count)
    pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
    return pixmap.tobytes("png")


def _pdf_candidates(
    document: NormalizedDocument, file_path: Path
) -> list[VisualCandidate]:
    blocks = {block.location: block for block in document.content}
    candidates: list[VisualCandidate] = []
    with fitz.open(file_path) as pdf:
        for index, page in enumerate(pdf):
            location = index + 1
            existing_text = blocks.get(location).text if location in blocks else ""
            sparse = len(existing_text.strip()) < max(
                0, config.VISION_MIN_TEXT_CHARS_PER_PAGE
            )
            visual_coverage = _page_image_coverage(page)
            image_heavy = visual_coverage >= max(
                0.0, min(1.0, config.VISION_MIN_IMAGE_COVERAGE)
            )
            if not sparse and not image_heavy:
                continue
            try:
                candidates.append(
                    VisualCandidate(
                        image_bytes=_render_pdf_page(page),
                        mime_type="image/png",
                        source_type=ContentBlockType.PAGE,
                        location=location,
                    )
                )
            except Exception:
                logger.exception(
                    "Failed to render visual candidate for %s page %s",
                    document.document_id,
                    location,
                )
                _append_once(
                    document.warnings,
                    f"Visual content could not be prepared for Page {location}.",
                )
    return candidates


def _pptx_candidates(
    document: NormalizedDocument, file_path: Path
) -> list[VisualCandidate]:
    candidates: list[VisualCandidate] = []
    presentation = Presentation(str(file_path))
    for index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
                width, height = image.size
            except Exception:
                logger.exception(
                    "Failed to read visual candidate for %s slide %s",
                    document.document_id,
                    index + 1,
                )
                _append_once(
                    document.warnings,
                    f"Visual content could not be prepared for Slide {index + 1}.",
                )
                continue
            if (
                width < max(1, config.VISION_MIN_EMBEDDED_IMAGE_WIDTH)
                or height < max(1, config.VISION_MIN_EMBEDDED_IMAGE_HEIGHT)
                or width * height > max(1, config.VISION_MAX_IMAGE_PIXELS)
            ):
                _append_once(
                    document.warnings,
                    f"An embedded visual on Slide {index + 1} was skipped by image safety limits.",
                )
                continue
            candidates.append(
                VisualCandidate(
                    image_bytes=image.blob,
                    mime_type=image.content_type,
                    source_type=ContentBlockType.SLIDE,
                    location=index + 1,
                )
            )
    return candidates


def _collect_candidates(
    document: NormalizedDocument, file_path: Path
) -> list[VisualCandidate]:
    if document.file_type == FileType.IMAGE:
        mime_type = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".webp": "image/webp",
        }.get(file_path.suffix.lower(), "application/octet-stream")
        return [
            VisualCandidate(
                image_bytes=file_path.read_bytes(),
                mime_type=mime_type,
                source_type=ContentBlockType.IMAGE,
                location=1,
            )
        ]
    if document.file_type == FileType.PDF:
        return _pdf_candidates(document, file_path)
    if document.file_type == FileType.PPTX:
        return _pptx_candidates(document, file_path)
    return []


def _merge_observation(
    document: NormalizedDocument,
    candidate: VisualCandidate,
    observation: VisualAnalysis,
) -> None:
    block = next(
        (
            item
            for item in document.content
            if item.type == candidate.source_type and item.location == candidate.location
        ),
        None,
    )
    if block is None:
        block = ContentBlock(
            type=candidate.source_type,
            location=candidate.location,
            text="",
        )
        document.content.append(block)

    evidence = observation.as_evidence_text(block.text)
    if not evidence:
        return
    block.text = f"{block.text.strip()}\n\n[Visual observations]\n{evidence}".strip()
    block.extra["visual_analysis"] = True
    block.extra["visual_source"] = "multimodal"


async def enrich_document_with_visuals(
    document: NormalizedDocument, file_path: Path
) -> NormalizedDocument:
    """Add visual evidence while preserving usable text on every failure."""
    if not config.MULTIMODAL_ENABLED:
        return document

    try:
        candidates = _collect_candidates(document, file_path)
    except Exception:  # document extraction already succeeded; visual add-on is optional
        logger.exception("Failed to identify visual candidates for %s", document.document_id)
        _append_once(document.warnings, "Visual content could not be prepared for analysis.")
        return document

    total_candidates = len(candidates)
    limit = max(1, config.VISION_MAX_ITEMS_PER_DOCUMENT)
    candidates = candidates[:limit]
    if not candidates:
        return document

    try:
        client = get_vision_client()
    except VisionError as exc:
        logger.warning(
            "Vision configuration unavailable for %s: %s",
            document.document_id,
            exc.message,
        )
        _append_once(document.warnings, "Visual analysis is currently unavailable.")
        document.metadata["multimodal"] = {
            "status": "unavailable",
            "attempted": 0,
            "succeeded": 0,
            "failed": len(candidates),
        }
        return document
    except Exception:
        logger.exception("Unexpected vision configuration failure for %s", document.document_id)
        _append_once(document.warnings, "Visual analysis is currently unavailable.")
        document.metadata["multimodal"] = {
            "status": "unavailable",
            "attempted": 0,
            "succeeded": 0,
            "failed": len(candidates),
        }
        return document

    cache: dict[str, VisualAnalysis | VisionError] = {}
    succeeded = 0
    failed = 0
    processed = 0
    document_timed_out = False
    loop = asyncio.get_running_loop()
    deadline = loop.time() + max(1.0, config.VISION_DOCUMENT_TIMEOUT_SECONDS)
    for candidate in candidates:
        remaining_seconds = deadline - loop.time()
        if remaining_seconds <= 0:
            document_timed_out = True
            break
        processed += 1
        digest = hashlib.sha256(candidate.image_bytes).hexdigest()
        cached = cache.get(digest)
        try:
            if isinstance(cached, VisionError):
                raise cached
            if cached is None:
                cached = await asyncio.wait_for(
                    client.analyze_image(
                        candidate.image_bytes,
                        candidate.mime_type,
                        candidate.source_type.value,
                        candidate.location,
                    ),
                    timeout=remaining_seconds,
                )
                cache[digest] = cached
            _merge_observation(document, candidate, cached)
            succeeded += 1
        except asyncio.TimeoutError:
            document_timed_out = True
            failed += 1
            _append_once(
                document.warnings,
                _safe_warning(candidate.source_type, candidate.location),
            )
            break
        except VisionError as exc:
            cache[digest] = exc
            failed += 1
            logger.warning(
                "Vision analysis failed for %s %s %s: %s",
                document.document_id,
                candidate.source_type.value,
                candidate.location,
                exc.message,
            )
            _append_once(
                document.warnings,
                _safe_warning(candidate.source_type, candidate.location),
            )
        except Exception:
            logger.exception(
                "Unexpected vision failure for %s %s %s",
                document.document_id,
                candidate.source_type.value,
                candidate.location,
            )
            cache[digest] = VisionResponseError("Unexpected vision processing failure.")
            failed += 1
            _append_once(
                document.warnings,
                _safe_warning(candidate.source_type, candidate.location),
            )

    status = "completed" if failed == 0 else "partial" if succeeded else "unavailable"
    document.metadata["multimodal"] = {
        "status": status,
        "provider": client.provider_name,
        "model": client.model_name,
        "attempted": processed,
        "succeeded": succeeded,
        "failed": failed,
        "skipped": max(0, total_candidates - processed),
    }
    if total_candidates > len(candidates):
        _append_once(
            document.warnings,
            f"Visual analysis was limited to {len(candidates)} items for this document.",
        )
    if document_timed_out:
        _append_once(
            document.warnings,
            "Visual analysis stopped after reaching the document time limit.",
        )
    if succeeded:
        document.status = DocumentStatus.PROCESSED
        document.warnings = [
            warning
            for warning in document.warnings
            if "later phase" not in warning.lower()
            and "requires configured" not in warning.lower()
        ]
    return document

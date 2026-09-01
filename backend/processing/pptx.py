"""PPTX processing using python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from processing.base import BaseProcessor
from processing.exceptions import DocumentProcessingError
from schemas.document import (
    ContentBlock,
    ContentBlockType,
    DocumentStatus,
    FileType,
    NormalizedDocument,
)


class PPTXProcessor(BaseProcessor):
    def process(
        self,
        file_path: Path,
        document_id: str,
        filename: str,
    ) -> NormalizedDocument:
        try:
            presentation = Presentation(str(file_path))
        except PackageNotFoundError as exc:
            raise DocumentProcessingError(f"Could not open PPTX: {exc}") from exc
        except Exception as exc:  # noqa: BLE001
            raise DocumentProcessingError(f"Could not open PPTX: {exc}") from exc

        content: list[ContentBlock] = []

        try:
            for index, slide in enumerate(presentation.slides):
                slide_number = index + 1
                title = self._extract_title(slide)
                text_parts = self._extract_text_boxes(slide)
                notes = self._extract_notes(slide)

                combined_text = "\n".join(part for part in text_parts if part)

                extra: dict = {}
                if title:
                    extra["title"] = title
                if notes:
                    extra["speaker_notes"] = notes

                content.append(
                    ContentBlock(
                        type=ContentBlockType.SLIDE,
                        location=slide_number,
                        text=combined_text,
                        extra=extra,
                    )
                )
        except Exception as exc:  # noqa: BLE001
            raise DocumentProcessingError(
                f"Failed to extract PPTX content: {exc}"
            ) from exc

        return NormalizedDocument(
            document_id=document_id,
            filename=filename,
            file_type=FileType.PPTX,
            status=DocumentStatus.PROCESSED,
            content=content,
            metadata={"slide_count": len(content)},
        )

    @staticmethod
    def _extract_title(slide) -> str | None:
        try:
            if slide.shapes.title and slide.shapes.title.text:
                return slide.shapes.title.text.strip()
        except Exception:  # noqa: BLE001
            pass
        return None

    @staticmethod
    def _extract_text_boxes(slide) -> list[str]:
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = "".join(run.text for run in paragraph.runs)
                if paragraph_text.strip():
                    parts.append(paragraph_text.strip())
        return parts

    @staticmethod
    def _extract_notes(slide) -> str | None:
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                return notes_text.strip() if notes_text and notes_text.strip() else None
        except Exception:  # noqa: BLE001
            pass
        return None
